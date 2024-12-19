import os
import json
import random
from pptx import Presentation
from transformers import pipeline

# Prompts to get variation in questions
prompt_variants = [
    "Generate a detailed question based on the context: {context}",
    "Generate a short question based on the context: {context}",
    "Generate a creative question based on the context: {context}",
    "Generate a factual question based on the context: {context}",
    "Generate a critical thinking question based on the context: {context}",
]

# Extract text from a .txt file and split into paragraphs.
def read_txt(file_path):
    
    with open(file_path, "r", encoding="utf-8") as file:
        text = file.read()
    paragraphs = [para.strip() for para in text.split("\n\n") if para.strip()]
    return paragraphs

# Extract text from a PowerPoint file and combine text for each slide.
def read_pptx(file_path):
    prs = Presentation(file_path)
    paragraphs = []
    for slide in prs.slides:
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                slide_text.append(shape.text)
        combined_text = " ".join(slide_text).strip()
        if combined_text:
            paragraphs.append(combined_text)
    return paragraphs

# Generate multiple diverse questions.
def generate_questions(context, num_questions):
    question_generator = pipeline("text2text-generation", model="google/flan-t5-base", device=0)
    questions = []
    seen_questions = set()
    
    for _ in range(num_questions * 2):  
        prompt_template = random.choice(prompt_variants)
        prompt = prompt_template.format(context=context)
        generated = question_generator(
            prompt,
            max_length=50,
            do_sample=True,
            top_k=50,
            top_p=0.9,
            temperature=0.8
        )[0]["generated_text"]

        # Add the question if it's unique
        if generated not in seen_questions:
            questions.append(generated)
            seen_questions.add(generated)

        if len(questions) >= num_questions:
            break
    
    return questions

# Process the input file and generate questions.
def process_file(input_path, num_questions, output_path):
    # Detect file type
    ext = os.path.splitext(input_path)[1].lower()
    if ext == ".txt":
        paragraphs = read_txt(input_path)
    elif ext == ".pptx":
        paragraphs = read_pptx(input_path)
    else:
        raise ValueError("Unsupported file format. Please provide a .txt or .pptx file.")

    results = []
    for i, paragraph in enumerate(paragraphs):
        questions = generate_questions(paragraph, num_questions)
        results.append({
            "id": str(i),
            "context": paragraph,
            "questions": questions,
        })

    with open(output_path, "w", encoding="utf-8") as f:
        json.dump(results, f, indent=4)
    print(f"Results saved to {output_path}")


def main():
    import argparse
    parser = argparse.ArgumentParser(description="Generate questions from input documents.")
    parser.add_argument("input_path", type=str, help="Path to the input file (TXT or PPTX).")
    parser.add_argument("--num_questions", type=int, default=1, help="Number of questions to generate per context.")
    parser.add_argument("--output", type=str, required=True, help="Path to save the output JSON.")
    args = parser.parse_args()

    process_file(args.input_path, args.num_questions, args.output)


if __name__ == "__main__":
    main()
